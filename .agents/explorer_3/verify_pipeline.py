import os
import sys
import io
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import pytesseract
from pytesseract import Output
from PIL import Image
import numpy as np
import cv2

# Set Tesseract executable path
pytesseract.pytesseract.tesseract_cmd = r'C:\Users\tvn183660\AppData\Local\Programs\Tesseract-OCR\tesseract.exe'

p1 = r'\\10.170.162.32\Data\00_KDTVN Common(KDTVN共通)\⑤Production Engineering(製造技術)\◆Iris EXP◆\Điều chỉnh\Athenal\Athena保証工程取り組み説明2025 VN.pptx'
p2 = r'\\10.170.162.32\Data\00_KDTVN Common(KDTVN共通)\⑤Production Engineering(製造技術)\◆Iris EXP◆\Điều chỉnh\Athenal\Athena保証工程　RaspberryPI問題点 VN.pptx'

def test_ocr_and_inpainting_on_sample():
    print("=== TESTING OCR & INPAINTING PIPELINE ===")
    prs = pptx.Presentation(p1)
    
    found_pic = None
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # check image size
                img_data = shape.image.blob
                img_pil = Image.open(io.BytesIO(img_data))
                w, h = img_pil.size
                if w > 300 and h > 150: # find a substantial diagram/screenshot
                    found_pic = (slide_idx, slide, shape, img_pil)
                    break
        if found_pic:
            break
            
    if not found_pic:
        print("No suitable picture found in File 1")
        return
        
    slide_idx, slide, shape, img_pil = found_pic
    print(f"Testing on Slide {slide_idx+1}, Shape: {shape.name}, Size: {img_pil.size}, Format: {img_pil.format}")
    
    # 1. Convert to OpenCV BGR format
    img_np = np.array(img_pil.convert('RGB'))
    img_bgr = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)
    
    # 2. Run Tesseract with data output
    data = pytesseract.image_to_data(img_pil, lang='jpn+eng', output_type=Output.DICT)
    
    detected_boxes = []
    n_boxes = len(data['text'])
    print(f"Total raw text items detected: {n_boxes}")
    for i in range(n_boxes):
        txt = data['text'][i].strip()
        conf = int(data['conf'][i])
        if conf > 30 and len(txt) > 0:
            x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
            detected_boxes.append((x, y, w, h, conf, txt))
            print(f"   [{conf}%] ({x}, {y}, {w}x{h}): {txt}")
            
    # 3. Test Inpainting on detected text boxes
    mask = np.zeros(img_bgr.shape[:2], dtype=np.uint8)
    for (x, y, w, h, conf, txt) in detected_boxes:
        pad = 2
        x0 = max(0, x - pad)
        y0 = max(0, y - pad)
        x1 = min(img_bgr.shape[1], x + w + pad)
        y1 = min(img_bgr.shape[0], y + h + pad)
        cv2.rectangle(mask, (x0, y0), (x1, y1), 255, -1)
        
    inpainted = cv2.inpaint(img_bgr, mask, inpaintRadius=3, flags=cv2.INPAINT_TELEA)
    print("Inpainting successful! Output shape:", inpainted.shape)
    
    # 4. Test Coordinate Transformation to Slide Space
    img_w, img_h = img_pil.size
    crop_l = getattr(shape, 'crop_left', 0.0) or 0.0
    crop_r = getattr(shape, 'crop_right', 0.0) or 0.0
    crop_t = getattr(shape, 'crop_top', 0.0) or 0.0
    crop_b = getattr(shape, 'crop_bottom', 0.0) or 0.0
    
    vis_x0 = crop_l * img_w
    vis_y0 = crop_t * img_h
    vis_w = (1.0 - crop_l - crop_r) * img_w
    vis_h = (1.0 - crop_t - crop_b) * img_h
    
    scale_x = shape.width / vis_w
    scale_y = shape.height / vis_h
    
    print(f"Shape bounds: left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
    print(f"Scale factors: scale_x={scale_x:.2f} EMUs/px, scale_y={scale_y:.2f} EMUs/px")
    
    for (x, y, w, h, conf, txt) in detected_boxes[:3]:
        slide_x = shape.left + int((x - vis_x0) * scale_x)
        slide_y = shape.top + int((y - vis_y0) * scale_y)
        slide_w = int(w * scale_x)
        slide_h = int(h * scale_y)
        print(f"Transformed box for '{txt}': slide_pos=({slide_x}, {slide_y}), slide_size=({slide_w}, {slide_h})")
        
    print("Verification completed successfully!")

if __name__ == '__main__':
    test_ocr_and_inpainting_on_sample()
