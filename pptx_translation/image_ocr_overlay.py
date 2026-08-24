"""
Image OCR, Inpainting, and Vietnamese Overlay Module.
Extracts embedded images from PPTX (including nested groups), runs Tesseract 5.5 OCR
to detect Japanese text, inpaints/erases original text, updates image blobs,
and creates coordinate-transformed Times New Roman Vietnamese text box overlays on slides.
"""

import os
import io
import re
from typing import List, Dict, Tuple, Optional
import numpy as np
import cv2
from PIL import Image
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
import pytesseract
from pytesseract import Output

from .translator_engine import PPTXTranslatorEngine
from .openxml_typography import apply_times_new_roman_complete

# Set Tesseract executable path on Windows
TESSERACT_EXE_CANDIDATES = [
    r"C:\Users\tvn183660\AppData\Local\Programs\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
]
for candidate in TESSERACT_EXE_CANDIDATES:
    if os.path.exists(candidate):
        pytesseract.pytesseract.tesseract_cmd = candidate
        break


class ImageOCROverlayProcessor:
    """
    Processes embedded PPTX images: OCR detection, adaptive inpainting,
    blob replacement, and slide coordinate mapped text box overlay.
    """

    def __init__(self, translator: Optional[PPTXTranslatorEngine] = None):
        self.translator = translator or PPTXTranslatorEngine()
        self.audit_log: List[Dict] = []

    def process_presentation_images(self, prs: pptx.Presentation) -> Dict[str, int]:
        """
        Scans all slides and groups, processes images containing Japanese text,
        erases Japanese text from image blobs, and inserts Vietnamese overlay text boxes.
        """
        stats = {
            "total_images_found": 0,
            "images_with_japanese": 0,
            "total_text_boxes_overlaid": 0,
        }

        for slide_idx, slide in enumerate(prs.slides, start=1):
            picture_items = list(self._find_all_pictures(slide.shapes))
            for shape, is_in_group, parent_group, (abs_left, abs_top) in picture_items:
                stats["total_images_found"] += 1
                try:
                    res = self._process_single_image(
                        slide=slide,
                        shape=shape,
                        slide_num=slide_idx,
                        abs_pos=(abs_left, abs_top),
                    )
                    if res["detected_count"] > 0:
                        stats["images_with_japanese"] += 1
                        stats["total_text_boxes_overlaid"] += res["detected_count"]
                except Exception as e:
                    self.audit_log.append({
                        "slide": slide_idx,
                        "shape_name": getattr(shape, "name", "unknown"),
                        "error": str(e),
                    })

        return stats

    def _find_all_pictures(
        self, shapes, parent_offset_x: int = 0, parent_offset_y: int = 0
    ):
        """
        Recursively yields (shape, is_in_group, parent_group, (abs_left, abs_top)) for all pictures,
        accumulating parent_offset_x and parent_offset_y across arbitrary nesting levels.
        """
        for shape in shapes:
            abs_left = parent_offset_x + getattr(shape, "left", 0)
            abs_top = parent_offset_y + getattr(shape, "top", 0)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                is_in_group = (parent_offset_x != 0 or parent_offset_y != 0)
                yield shape, is_in_group, None, (abs_left, abs_top)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._find_all_pictures(
                    shape.shapes,
                    parent_offset_x=abs_left,
                    parent_offset_y=abs_top,
                )

    def _process_single_image(
        self, slide, shape, slide_num: int, abs_pos: Tuple[int, int]
    ) -> Dict:
        """Processes one image: OCR -> Inpaint -> Blob replace -> Slide text box overlay."""
        result = {"detected_count": 0, "boxes": []}
        
        # 1. Load image from blob
        try:
            image_blob = shape.image.blob
            img_pil = Image.open(io.BytesIO(image_blob))
        except Exception:
            return result

        orig_w, orig_h = img_pil.size
        if orig_w < 80 or orig_h < 40:
            return result

        # 2. Preprocess image for OCR
        img_np = np.array(img_pil.convert("RGB"))
        h_orig, w_orig = img_np.shape[:2]
        
        # Scale 2x if small for better Kanji OCR
        scale = 2.0 if max(h_orig, w_orig) < 1200 else 1.0
        if scale > 1.0:
            resized_np = cv2.resize(
                img_np, (int(w_orig * scale), int(h_orig * scale)), interpolation=cv2.INTER_CUBIC
            )
        else:
            resized_np = img_np

        gray = cv2.cvtColor(resized_np, cv2.COLOR_RGB2GRAY)
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced_gray = clahe.apply(gray)
        enhanced_pil = Image.fromarray(enhanced_gray)

        # 3. Run OCR with Tesseract
        try:
            ocr_data = pytesseract.image_to_data(
                enhanced_pil, lang="jpn+eng", config="--psm 11 --oem 1", output_type=Output.DICT
            )
        except Exception:
            return result

        # 4. Filter and Cluster Bounding Boxes
        raw_boxes = []
        n_items = len(ocr_data["text"])
        for i in range(n_items):
            txt = ocr_data["text"][i].strip()
            conf = int(ocr_data["conf"][i]) if ocr_data["conf"][i] != "-1" else 0
            if conf >= 30 and txt and self.translator.contains_japanese(txt):
                x = int(ocr_data["left"][i] / scale)
                y = int(ocr_data["top"][i] / scale)
                w = int(ocr_data["width"][i] / scale)
                h = int(ocr_data["height"][i] / scale)
                if w > 4 and h > 4 and (w / h < 30) and (h / w < 30):
                    raw_boxes.append({"x": x, "y": y, "w": w, "h": h, "text": txt, "conf": conf})

        if not raw_boxes:
            return result

        clustered_boxes = self._cluster_bounding_boxes(raw_boxes)
        if not clustered_boxes:
            return result

        # 5. Inpaint / Erase Japanese Text on OpenCV BGR Image
        img_bgr = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)
        for box in clustered_boxes:
            x, y, w, h = box["x"], box["y"], box["w"], box["h"]
            pad = 2
            x0 = max(0, x - pad)
            y0 = max(0, y - pad)
            x1 = min(w_orig, x + w + pad)
            y1 = min(h_orig, y + h + pad)

            # Sample border pixels to compute variance
            border_mask = np.zeros(img_bgr.shape[:2], dtype=bool)
            border_mask[max(0, y0 - 2):min(h_orig, y1 + 2), max(0, x0 - 2):min(w_orig, x1 + 2)] = True
            border_mask[y0:y1, x0:x1] = False
            
            border_pixels = img_bgr[border_mask]
            if len(border_pixels) > 0:
                std_dev = np.std(border_pixels)
                if std_dev < 12.0:
                    # Flat / uniform color background -> clean median fill
                    median_color = np.median(border_pixels, axis=0).astype(int).tolist()
                    cv2.rectangle(img_bgr, (x0, y0), (x1, y1), median_color, -1)
                else:
                    # Textured background -> OpenCV Telea inpaint
                    mask = np.zeros(img_bgr.shape[:2], dtype=np.uint8)
                    mask[y0:y1, x0:x1] = 255
                    img_bgr = cv2.inpaint(img_bgr, mask, inpaintRadius=3, flags=cv2.INPAINT_TELEA)

        # 6. Update Image Blob in PPTX Shape Relationship
        self._replace_shape_blob(shape, img_bgr)

        # 7. Translate and Add Overlay Text Boxes to Slide
        img_w, img_h = orig_w, orig_h
        crop_l = getattr(shape, "crop_left", 0.0) or 0.0
        crop_r = getattr(shape, "crop_right", 0.0) or 0.0
        crop_t = getattr(shape, "crop_top", 0.0) or 0.0
        crop_b = getattr(shape, "crop_bottom", 0.0) or 0.0

        vis_x0 = crop_l * img_w
        vis_y0 = crop_t * img_h
        vis_w = max(1.0, (1.0 - crop_l - crop_r) * img_w)
        vis_h = max(1.0, (1.0 - crop_t - crop_b) * img_h)

        abs_left, abs_top = abs_pos
        scale_x = shape.width / vis_w
        scale_y = shape.height / vis_h

        for box in clustered_boxes:
            ja_text = box["text"]
            vn_text = self.translator.translate_text(ja_text)
            
            x, y, w, h = box["x"], box["y"], box["w"], box["h"]
            slide_x = abs_left + int((x - vis_x0) * scale_x)
            slide_y = abs_top + int((y - vis_y0) * scale_y)
            # Expansion margin for Vietnamese text length
            slide_w = max(Emu(180000), int(w * scale_x * 1.25))
            slide_h = max(Emu(120000), int(h * scale_y * 1.15))

            self._create_overlay_textbox(slide, slide_x, slide_y, slide_w, slide_h, vn_text)
            
            result["detected_count"] += 1
            result["boxes"].append({
                "ja_text": ja_text,
                "vn_text": vn_text,
                "box": (slide_x, slide_y, slide_w, slide_h),
            })

        return result

    def _cluster_bounding_boxes(self, raw_boxes: List[Dict]) -> List[Dict]:
        """Clusters adjacent horizontal character bounding boxes into coherent phrase boxes."""
        if not raw_boxes:
            return []

        clusters = [dict(b) for b in raw_boxes]

        while True:
            best_pair = None
            best_gap = float("inf")

            for i in range(len(clusters)):
                for j in range(i + 1, len(clusters)):
                    c1 = clusters[i]
                    c2 = clusters[j]
                    max_h = max(c1["h"], c2["h"])
                    vert_diff = abs(c1["y"] - c2["y"])

                    if vert_diff <= max_h * 0.7:
                        left_c, right_c = (c1, c2) if c1["x"] <= c2["x"] else (c2, c1)
                        gap = right_c["x"] - (left_c["x"] + left_c["w"])
                        if -max_h * 0.5 <= gap <= max_h * 1.5:
                            abs_gap = abs(gap)
                            if abs_gap < best_gap:
                                best_gap = abs_gap
                                best_pair = (i, j, left_c, right_c)

            if best_pair is not None:
                i, j, left_c, right_c = best_pair
                merged_box = {
                    "x": min(left_c["x"], right_c["x"]),
                    "y": min(left_c["y"], right_c["y"]),
                    "w": max(left_c["x"] + left_c["w"], right_c["x"] + right_c["w"]) - min(left_c["x"], right_c["x"]),
                    "h": max(left_c["y"] + left_c["h"], right_c["y"] + right_c["y"]) - min(left_c["y"], right_c["y"]),
                    "text": left_c["text"] + right_c["text"],
                    "conf": min(left_c["conf"], right_c["conf"]),
                }
                clusters.pop(max(i, j))
                clusters.pop(min(i, j))
                clusters.append(merged_box)
            else:
                break

        clusters.sort(key=lambda b: (b["y"], b["x"]))
        return clusters

    def _replace_shape_blob(self, shape, img_bgr: np.ndarray) -> None:
        """Encodes OpenCV image and replaces the PPTX ImagePart blob."""
        success, buffer = cv2.imencode(".png", img_bgr)
        if not success:
            return
        new_blob = buffer.tobytes()

        try:
            blip_elem = shape._element.blipFill.blip
            embed_rId = blip_elem.attrib[
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            ]
            image_part = shape.part.related_parts[embed_rId]
            image_part._blob = new_blob
        except Exception:
            pass

    def _create_overlay_textbox(
        self, slide, slide_x: int, slide_y: int, slide_w: int, slide_h: int, vn_text: str
    ) -> None:
        """Creates a styled, transparent, Times New Roman overlay textbox on slide."""
        try:
            tx_box = slide.shapes.add_textbox(slide_x, slide_y, slide_w, slide_h)
            tf = tx_box.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)

            p = tf.paragraphs[0]
            p.text = vn_text
            p.alignment = PP_ALIGN.LEFT

            # Calculate optimal font point size matching original box height
            h_pt = slide_h / 12700.0
            font_size_pt = max(7.5, min(22.0, h_pt * 0.72))
            p.font.size = Pt(font_size_pt)
            p.font.name = "Times New Roman"
            p.font.color.rgb = RGBColor(25, 25, 25)

            # Apply OpenXML font normalization
            apply_times_new_roman_complete(p)
            for r in p.runs:
                apply_times_new_roman_complete(r)
        except Exception:
            pass
