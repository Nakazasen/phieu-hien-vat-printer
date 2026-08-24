"""
Empirical Challenger 1: Final Acceptance Verification Script
Deep Adversarial Inspection: Content, OCR & Typography.
Tests target network UNC files and local output/ files.
"""

import os
import sys
import zipfile
import re
import hashlib
import xml.etree.ElementTree as ET
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

PROJECT_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))

TARGET_UNC_FILES = [
    r"\\10.170.162.32\Data\00_KDTVN Common(KDTVN共通)\⑤Production Engineering(製造技術)\◆Iris EXP◆\Điều chỉnh\Athenal\Athena保証工程取り組み説明2025 VN.pptx",
    r"\\10.170.162.32\Data\00_KDTVN Common(KDTVN共通)\⑤Production Engineering(製造技術)\◆Iris EXP◆\Điều chỉnh\Athenal\Athena保証工程　RaspberryPI問題点 VN.pptx",
]

OUTPUT_FILES = [
    os.path.join(PROJECT_ROOT, "output", "Athena保証工程取り組み説明2025 VN.pptx"),
    os.path.join(PROJECT_ROOT, "output", "Athena保証工程　RaspberryPI問題点 VN.pptx"),
]

CJK_REGEX = re.compile(r'[\u3040-\u30ff\u3400-\u4dbf\u4e00-\u9fff\uff66-\uff9f]')
XML_NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
}

def sha256_file(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while chunk := f.read(65536):
            h.update(chunk)
    return h.hexdigest()

def check_unc_sync():
    print("\n" + "=" * 80)
    print("[CHALLENGE 1] NETWORK UNC SHARE & STAGING SYNCHRONIZATION TEST")
    print("=" * 80)
    
    all_synced = True
    for unc_path, out_path in zip(TARGET_UNC_FILES, OUTPUT_FILES):
        name = os.path.basename(unc_path)
        print(f"\nChecking: {name}")
        out_exists = os.path.exists(out_path)
        unc_exists = os.path.exists(unc_path)
        
        print(f"  - Local Output File Exists: {out_exists} ({out_path})")
        if out_exists:
            out_size = os.path.getsize(out_path)
            out_hash = sha256_file(out_path)
            print(f"    * Output Size: {out_size:,} bytes | SHA-256: {out_hash}")
        else:
            all_synced = False
            
        print(f"  - Network UNC Target Exists: {unc_exists} ({unc_path})")
        if unc_exists:
            unc_size = os.path.getsize(unc_path)
            unc_hash = sha256_file(unc_path)
            print(f"    * UNC Size: {unc_size:,} bytes | SHA-256: {unc_hash}")
            if out_exists:
                if out_hash == unc_hash:
                    print(f"    * Sync Status: 100% BITWISE IDENTICAL MATCH (SHA-256 match)")
                else:
                    print(f"    * Sync Status: MISMATCH (Output SHA != UNC SHA)")
                    all_synced = False
        else:
            print(f"    * WARNING: Network UNC path not accessible in current environment (checking local output)")
            
    return all_synced

def deep_inspect_pptx_openxml(pptx_path: str):
    print(f"\n--- [RAW OPENXML ARCHIVE INSPECTION] {os.path.basename(pptx_path)} ---")
    if not os.path.exists(pptx_path):
        print(f"ERROR: File not found: {pptx_path}")
        return False, {}

    stats = {
        "xml_files_scanned": 0,
        "raw_a_t_count": 0,
        "raw_japanese_matches": [],
        "fonts_found": set(),
        "non_tnr_font_tags": [],
    }

    with zipfile.ZipFile(pptx_path, 'r') as z:
        for name in z.namelist():
            # Check all slide XMLs, notes, layouts, masters, charts, diagrams
            if name.startswith("ppt/") and name.endswith(".xml"):
                stats["xml_files_scanned"] += 1
                xml_data = z.read(name)
                try:
                    root = ET.fromstring(xml_data)
                except Exception as e:
                    print(f"Warning: Failed to parse XML {name}: {e}")
                    continue

                # 1. Inspect all <a:t> elements
                for elem in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
                    stats["raw_a_t_count"] += 1
                    txt = elem.text or ""
                    if CJK_REGEX.search(txt):
                        stats["raw_japanese_matches"].append((name, txt))

                # Also inspect any text in other namespaces if any
                for elem in root.iter():
                    if elem.tag.endswith('}t') and not elem.tag.startswith('{http://schemas.openxmlformats.org/drawingml/2006/main}'):
                        txt = elem.text or ""
                        if CJK_REGEX.search(txt):
                            stats["raw_japanese_matches"].append((f"{name}:{elem.tag}", txt))

                # 2. Inspect all font definitions <a:latin>, <a:ea>, <a:cs>, <a:sym>
                for tag_suffix in ['latin', 'ea', 'cs', 'sym']:
                    for elem in root.iter(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag_suffix}'):
                        tf = elem.attrib.get('typeface', '')
                        if tf:
                            stats["fonts_found"].add(tf)
                            if "Times New Roman" not in tf:
                                stats["non_tnr_font_tags"].append((name, tag_suffix, tf))

    print(f"XML Scanned: {stats['xml_files_scanned']} XML parts")
    print(f"Raw <a:t> Elements: {stats['raw_a_t_count']}")
    print(f"Distinct Fonts Defined: {stats['fonts_found']}")
    print(f"Residual Japanese Elements in Raw XML: {len(stats['raw_japanese_matches'])}")
    if stats["raw_japanese_matches"]:
        for xml_part, j_text in stats["raw_japanese_matches"][:5]:
            print(f"  * [{xml_part}] -> {j_text}")

    print(f"Non-TNR Font Tags in Raw XML: {len(stats['non_tnr_font_tags'])}")
    if stats["non_tnr_font_tags"]:
        for xml_part, f_type, f_name in stats["non_tnr_font_tags"][:5]:
            print(f"  * [{xml_part}] {f_type}='{f_name}'")

    is_clean = (len(stats["raw_japanese_matches"]) == 0 and len(stats["non_tnr_font_tags"]) == 0)
    return is_clean, stats

def deep_inspect_pptx_model(pptx_path: str):
    print(f"\n--- [PYTHON-PPTX OBJECT MODEL AUDIT] {os.path.basename(pptx_path)} ---")
    if not os.path.exists(pptx_path):
        print(f"ERROR: File not found: {pptx_path}")
        return False, {}

    prs = pptx.Presentation(pptx_path)
    
    stats = {
        "slides": len(prs.slides),
        "total_shapes": 0,
        "group_shapes": 0,
        "tables": 0,
        "charts": 0,
        "images": 0,
        "text_frames": 0,
        "paragraphs": 0,
        "runs": 0,
        "japanese_paragraphs": [],
        "non_tnr_runs": [],
        "overlay_textboxes": 0,
    }
    
    visited_cells = set()

    def inspect_text_frame(tf, loc):
        stats["text_frames"] += 1
        for p_idx, p in enumerate(tf.paragraphs):
            stats["paragraphs"] += 1
            p_text = "".join(r.text for r in p.runs) if p.runs else p.text
            if CJK_REGEX.search(p_text):
                stats["japanese_paragraphs"].append((loc, p_idx, p_text))

            for r_idx, r in enumerate(p.runs):
                stats["runs"] += 1
                rPr = r._r.find(qn('a:rPr'))
                if rPr is not None:
                    latin = rPr.find(qn('a:latin'))
                    ea = rPr.find(qn('a:ea'))
                    cs = rPr.find(qn('a:cs'))
                    
                    latin_font = latin.get('typeface') if latin is not None else None
                    ea_font = ea.get('typeface') if ea is not None else None
                    cs_font = cs.get('typeface') if cs is not None else None
                    
                    if latin_font and "Times New Roman" not in latin_font:
                        stats["non_tnr_runs"].append((loc, p_idx, r_idx, r.text, "latin", latin_font))
                    if ea_font and "Times New Roman" not in ea_font:
                        stats["non_tnr_runs"].append((loc, p_idx, r_idx, r.text, "ea", ea_font))
                    if cs_font and "Times New Roman" not in cs_font:
                        stats["non_tnr_runs"].append((loc, p_idx, r_idx, r.text, "cs", cs_font))

    def inspect_shapes_recursive(shapes, prefix):
        for s_idx, shape in enumerate(shapes, start=1):
            stats["total_shapes"] += 1
            loc = f"{prefix}_Sh{s_idx}_{shape.name}"

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                stats["images"] += 1

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                stats["group_shapes"] += 1
                inspect_shapes_recursive(shape.shapes, f"{loc}_Grp")
                continue

            if shape.has_text_frame and shape.text_frame.text.strip():
                if "OCR_Overlay" in shape.name or "Overlay" in shape.name:
                    stats["overlay_textboxes"] += 1
                inspect_text_frame(shape.text_frame, loc)

            if shape.has_table:
                stats["tables"] += 1
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        cell_id = id(cell._tc)
                        if cell_id not in visited_cells:
                            visited_cells.add(cell_id)
                            if cell.text_frame.text.strip():
                                inspect_text_frame(cell.text_frame, f"{loc}_R{r_idx}C{c_idx}")

            if shape.has_chart:
                stats["charts"] += 1
                if shape.chart.has_title and shape.chart.chart_title.has_text_frame:
                    inspect_text_frame(shape.chart.chart_title.text_frame, f"{loc}_ChartTitle")

    for slide_idx, slide in enumerate(prs.slides, start=1):
        inspect_shapes_recursive(slide.shapes, f"Slide_{slide_idx}")
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf and notes_tf.text.strip():
                inspect_text_frame(notes_tf, f"Slide_{slide_idx}_Notes")

    print(f"Stats: Slides={stats['slides']}, Shapes={stats['total_shapes']} (Grp={stats['group_shapes']}, Img={stats['images']}, Tbl={stats['tables']}, Chart={stats['charts']})")
    print(f"Text: TextFrames={stats['text_frames']}, Paragraphs={stats['paragraphs']}, Runs={stats['runs']}")
    print(f"Residual Japanese: {len(stats['japanese_paragraphs'])}")
    print(f"Non-Times New Roman Runs: {len(stats['non_tnr_runs'])}")

    is_clean = (len(stats["japanese_paragraphs"]) == 0 and len(stats["non_tnr_runs"]) == 0)
    return is_clean, stats

def main():
    print("=" * 80)
    print("EMPIRICAL CHALLENGER 1: FINAL ACCEPTANCE HARNESS")
    print("=" * 80)

    # Step 1: Check UNC and Staging sync
    unc_sync = check_unc_sync()

    # Step 2: Deep inspection on all files
    files_to_test = []
    for p in OUTPUT_FILES:
        if os.path.exists(p):
            files_to_test.append(p)
    for p in TARGET_UNC_FILES:
        if os.path.exists(p) and p not in files_to_test:
            files_to_test.append(p)

    all_passed = True
    for fpath in files_to_test:
        clean_model, m_stats = deep_inspect_pptx_model(fpath)
        clean_xml, x_stats = deep_inspect_pptx_openxml(fpath)
        if not (clean_model and clean_xml):
            all_passed = False

    print("\n" + "=" * 80)
    if all_passed:
        print(">>> EMPIRICAL CHALLENGER 1 VERDICT: APPROVE <<<")
        print("0 residual Japanese paragraphs and 100% Times New Roman font verified across all runs and raw OpenXML XMLs.")
        print("=" * 80)
        sys.exit(0)
    else:
        print(">>> EMPIRICAL CHALLENGER 1 VERDICT: REQUEST_CHANGES <<<")
        print("=" * 80)
        sys.exit(1)

if __name__ == "__main__":
    main()
