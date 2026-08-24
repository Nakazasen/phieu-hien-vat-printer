"""
OpenXML Typography & Layout Normalizer for PPTX Presentations.
Strictly enforces Times New Roman font across Latin, East Asian (EA), and Complex Script (CS)
typefaces in PresentationML DrawingML nodes, configures autofit, and prevents text overflow.
"""

from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn


def apply_times_new_roman_complete(element) -> None:
    """
    Directly modifies the OpenXML <a:rPr>, <a:defRPr>, or <a:endParaRPr> node
    to enforce Times New Roman across latin, ea (East Asian), and cs (Complex Script),
    and sets language to Vietnamese (vi-VN).
    """
    if element is None:
        return

    # Handle python-pptx Run object
    if hasattr(element, '_r'):
        rPr = element._r.get_or_add_rPr()
        _set_font_nodes(rPr)
    # Handle python-pptx Paragraph object
    elif hasattr(element, '_p'):
        p = element._p
        pPr = p.get_or_add_pPr()
        
        # 1. Update defRPr in pPr
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is None:
            defRPr = OxmlElement('a:defRPr')
            pPr.append(defRPr)
        _set_font_nodes(defRPr)
        
        # 2. Update endParaRPr
        endParaRPr = p.find(qn('a:endParaRPr'))
        if endParaRPr is None:
            endParaRPr = OxmlElement('a:endParaRPr')
            p.append(endParaRPr)
        _set_font_nodes(endParaRPr)
        return
    else:
        _set_font_nodes(element)


def _set_font_nodes(rPr_node) -> None:
    """Sets latin, ea, and cs font family tags to Times New Roman and lang to vi-VN."""
    if rPr_node is None:
        return

    # 1. Latin Typeface
    latin = rPr_node.find(qn('a:latin'))
    if latin is None:
        latin = OxmlElement('a:latin')
        rPr_node.append(latin)
    latin.set('typeface', 'Times New Roman')
    latin.set('pitchFamily', '18')
    latin.set('charset', '0')

    # 2. East Asian Typeface (Overwrites MS Gothic, Meiryo, Yu Gothic)
    ea = rPr_node.find(qn('a:ea'))
    if ea is None:
        ea = OxmlElement('a:ea')
        rPr_node.append(ea)
    ea.set('typeface', 'Times New Roman')
    ea.set('pitchFamily', '18')
    ea.set('charset', '0')

    # 3. Complex Script Typeface
    cs = rPr_node.find(qn('a:cs'))
    if cs is None:
        cs = OxmlElement('a:cs')
        rPr_node.append(cs)
    cs.set('typeface', 'Times New Roman')
    cs.set('pitchFamily', '18')
    cs.set('charset', '0')

    # 4. Language tag
    rPr_node.set('lang', 'vi-VN')


class OpenXMLTypographyNormalizer:
    """
    Normalizes typography, margins, word-wrap, and autofit across all
    TextFrames, Paragraphs, and Runs in a Presentation.
    """

    @classmethod
    def normalize_text_frame(cls, text_frame, expansion_ratio: float = 1.0) -> None:
        """
        Applies full Times New Roman font normalization and layout optimization
        to a python-pptx TextFrame.
        """
        if text_frame is None:
            return

        # 1. Enable Word Wrap
        text_frame.word_wrap = True

        # 2. Compress internal padding margins to prevent text overflow
        try:
            text_frame.margin_left = Inches(0.03)
            text_frame.margin_right = Inches(0.03)
            text_frame.margin_top = Inches(0.02)
            text_frame.margin_bottom = Inches(0.02)
        except Exception:
            pass

        # 3. Enforce Native OpenXML Normal Autofit (<a:normAutofit/>)
        try:
            bodyPr = text_frame._txBody.bodyPr
            for tag_name in ['a:noAutofit', 'a:spAutoFit']:
                elem = bodyPr.find(qn(tag_name))
                if elem is not None:
                    bodyPr.remove(elem)
            
            if bodyPr.find(qn('a:normAutofit')) is None:
                bodyPr.append(OxmlElement('a:normAutofit'))
        except Exception:
            pass

        # 4. Apply Times New Roman across all paragraphs and runs
        for paragraph in text_frame.paragraphs:
            apply_times_new_roman_complete(paragraph)
            
            for run in paragraph.runs:
                apply_times_new_roman_complete(run)
                # Also set standard python-pptx property as fallback
                try:
                    run.font.name = "Times New Roman"
                except Exception:
                    pass

                # Dynamic font scaling if text expansion is large (> 1.45x)
                if expansion_ratio > 1.45 and run.font.size and run.font.size.pt > 11:
                    scale_factor = 0.88 if expansion_ratio <= 1.7 else 0.80
                    new_pt = max(8.0, round(run.font.size.pt * scale_factor, 1))
                    run.font.size = Pt(new_pt)
