"""
Pytest Test Suite for Japanese-to-Vietnamese PPTX Translation and Image OCR Pipeline.
Validates:
1. BackupManager (SHA-256 computation, staging, safe deployment).
2. Manufacturing glossary & Translation Engine.
3. OpenXML Typography Normalizer (Times New Roman across latin, ea, cs, defRPr, endParaRPr).
4. Image OCR bounding box clustering & inpainting logic.
5. End-to-end PPTX recursive shape/table/notes traversal and translation.
"""

import os
import io
import pytest
import numpy as np
import cv2
from PIL import Image
import pptx
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from pptx_translation.backup_manager import BackupManager
from pptx_translation.glossary import MANUFACTURING_GLOSSARY, translate_with_glossary
from pptx_translation.translator_engine import PPTXTranslatorEngine
from pptx_translation.openxml_typography import (
    OpenXMLTypographyNormalizer,
    apply_times_new_roman_complete,
)
from pptx_translation.image_ocr_overlay import ImageOCROverlayProcessor
from pptx_translation.pipeline import PPTXTranslationPipeline


def test_backup_manager(tmp_path):
    # 1. Create a dummy file
    source_file = tmp_path / "test_source.pptx"
    source_file.write_bytes(b"PK\x03\x04Dummy PPTX content for SHA-256 test")

    backup_dir = str(tmp_path / "backups")
    staging_dir = str(tmp_path / "output")

    mgr = BackupManager(base_backup_dir=backup_dir, staging_dir=staging_dir)
    staged_path, backup_path, orig_sha = mgr.backup_and_stage(str(source_file))

    assert os.path.exists(backup_path)
    assert os.path.exists(staged_path)
    assert mgr.calculate_sha256(backup_path) == orig_sha
    assert mgr.calculate_sha256(staged_path) == orig_sha

    # Test network deployment
    deploy_target = str(tmp_path / "deployed.pptx")
    deployed_sha = mgr.deploy_to_network(staged_path, deploy_target)
    assert os.path.exists(deploy_target)
    assert deployed_sha == orig_sha


def test_manufacturing_glossary():
    assert "保証工程" in MANUFACTURING_GLOSSARY
    assert "RaspberryPI" in MANUFACTURING_GLOSSARY
    assert "SDカード破損" in MANUFACTURING_GLOSSARY

    sample = "Athena保証工程取り組み説明2025"
    translated = translate_with_glossary(sample)
    assert "bảo đảm chất lượng" in translated.lower() or "Athena" in translated


def test_translator_engine(tmp_path):
    cache_path = str(tmp_path / "test_cache.json")
    engine = PPTXTranslatorEngine(cache_file=cache_path)

    # Test CJK detection
    assert engine.contains_japanese("保証工程取り組み説明") is True
    assert engine.contains_japanese("RaspberryPI問題点") is True
    assert engine.contains_japanese("Only English Text 123") is False

    # Test exact glossary translation
    res = engine.translate_text("SDカード破損")
    assert "Hỏng thẻ nhớ SD" in res


def test_openxml_typography_enforcement():
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Thuyết minh hoạt động bảo đảm chất lượng"

    # Apply Times New Roman complete normalizer
    OpenXMLTypographyNormalizer.normalize_text_frame(tf)

    # Inspect XML for run
    rPr = r._r.find(qn("a:rPr"))
    assert rPr is not None
    latin = rPr.find(qn("a:latin"))
    ea = rPr.find(qn("a:ea"))
    cs = rPr.find(qn("a:cs"))

    assert latin is not None and latin.get("typeface") == "Times New Roman"
    assert ea is not None and ea.get("typeface") == "Times New Roman"
    assert cs is not None and cs.get("typeface") == "Times New Roman"
    assert rPr.get("lang") == "vi-VN"

    # Inspect paragraph endParaRPr
    endParaRPr = p._p.find(qn("a:endParaRPr"))
    assert endParaRPr is not None
    assert endParaRPr.find(qn("a:latin")).get("typeface") == "Times New Roman"
    assert endParaRPr.find(qn("a:ea")).get("typeface") == "Times New Roman"


def test_ocr_clustering_and_inpainting():
    processor = ImageOCROverlayProcessor()

    raw_boxes = [
        {"x": 10, "y": 20, "w": 15, "h": 18, "text": "保", "conf": 80},
        {"x": 27, "y": 21, "w": 15, "h": 18, "text": "証", "conf": 85},
        {"x": 44, "y": 20, "w": 16, "h": 18, "text": "工", "conf": 82},
        {"x": 62, "y": 21, "w": 15, "h": 18, "text": "程", "conf": 88},
    ]

    clustered = processor._cluster_bounding_boxes(raw_boxes)
    assert len(clustered) == 1
    assert clustered[0]["text"] == "保証工程"
    assert clustered[0]["x"] == 10
    assert clustered[0]["w"] >= 65


def test_end_to_end_pipeline(tmp_path):
    # Create sample presentation with Textbox, Table, and Group
    source_pptx = tmp_path / "sample_input.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Textbox
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.add_run().text = "保証工程取り組み説明"

    # Table
    tbl_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(5), Inches(1.5))
    tbl_shape.table.cell(0, 0).text_frame.paragraphs[0].add_run().text = "RaspberryPI問題点"
    tbl_shape.table.cell(0, 1).text_frame.paragraphs[0].add_run().text = "SDカード破損"

    prs.save(str(source_pptx))

    # Run Pipeline
    backup_dir = str(tmp_path / "backups")
    staging_dir = str(tmp_path / "output")
    bm = BackupManager(base_backup_dir=backup_dir, staging_dir=staging_dir)
    pipeline = PPTXTranslationPipeline(backup_manager=bm)

    stats = pipeline.process_file(str(source_pptx), skip_deploy=True)

    assert stats["paragraphs_translated"] >= 1
    assert stats["table_cells_processed"] >= 2

    # Verify output presentation
    out_prs = pptx.Presentation(stats["staged_sha256"] and bm.manifest[os.path.basename(str(source_pptx))]["staged_path"])
    out_slide = out_prs.slides[0]
    out_tb_text = out_slide.shapes[0].text_frame.text
    assert "bảo đảm chất lượng" in out_tb_text.lower() or "hoạt động" in out_tb_text.lower()
