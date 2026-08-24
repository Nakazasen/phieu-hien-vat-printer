"""
End-to-End PPTX Translation and OCR Pipeline Orchestrator.
Coordinates backup staging, recursive shape & table translation, OpenXML typography
normalization, embedded image OCR inpainting & overlay, and verified network deployment.
"""

import os
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Set, Optional

from .backup_manager import BackupManager
from .translator_engine import PPTXTranslatorEngine
from .openxml_typography import OpenXMLTypographyNormalizer, apply_times_new_roman_complete
from .image_ocr_overlay import ImageOCROverlayProcessor


class PPTXTranslationPipeline:
    """
    Executes the complete translation, typography normalization,
    and image OCR overlay workflow for PowerPoint presentations.
    """

    def __init__(
        self,
        backup_manager: Optional[BackupManager] = None,
        translator: Optional[PPTXTranslatorEngine] = None,
    ):
        self.backup_mgr = backup_manager or BackupManager()
        self.translator = translator or PPTXTranslatorEngine()
        self.ocr_processor = ImageOCROverlayProcessor(translator=self.translator)
        self.visited_cells: Set[int] = set()

    def process_file(
        self, source_path: str, target_unc_path: Optional[str] = None, skip_deploy: bool = False
    ) -> Dict:
        """
        Executes the end-to-end pipeline for a single PPTX presentation.
        """
        deploy_target = target_unc_path or source_path
        
        # 1. Backup & Stage
        staged_path, backup_path, orig_sha = self.backup_mgr.backup_and_stage(source_path)
        
        # 2. Open Presentation
        prs = pptx.Presentation(staged_path)
        self.visited_cells.clear()

        stats = {
            "file_name": os.path.basename(source_path),
            "total_slides": len(prs.slides),
            "paragraphs_translated": 0,
            "table_cells_processed": 0,
            "slide_notes_processed": 0,
            "images_found": 0,
            "images_with_japanese": 0,
            "image_overlay_boxes": 0,
            "original_sha256": orig_sha,
            "backup_path": backup_path,
        }

        # 3. Recursive Traversal & Translation of Native Text Elements
        for slide_idx, slide in enumerate(prs.slides, start=1):
            # A. Slide Shapes (including recursive GroupShapes and Tables)
            self._process_shapes_collection(slide.shapes, stats)

            # B. Slide Notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf and notes_tf.text.strip():
                    self._translate_text_frame(notes_tf, stats)
                    stats["slide_notes_processed"] += 1

        # 4. Embedded Image OCR, Inpainting, and Vietnamese Overlay
        image_stats = self.ocr_processor.process_presentation_images(prs)
        stats["images_found"] = image_stats["total_images_found"]
        stats["images_with_japanese"] = image_stats["images_with_japanese"]
        stats["image_overlay_boxes"] = image_stats["total_text_boxes_overlaid"]

        # 5. Save Processed Staging File
        prs.save(staged_path)
        stats["staged_sha256"] = self.backup_mgr.calculate_sha256(staged_path)

        # 6. Deploy to Network Share (if not skipped)
        if not skip_deploy:
            deployed_sha = self.backup_mgr.deploy_to_network(staged_path, deploy_target)
            stats["deployed_sha256"] = deployed_sha
            stats["deploy_target"] = deploy_target

        return stats

    def _process_shapes_collection(self, shapes, stats: Dict) -> None:
        """Recursively processes shapes, nested GroupShapes, Tables, and Charts."""
        for shape in shapes:
            # Group Shape (Recursive descent)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._process_shapes_collection(shape.shapes, stats)
                continue

            # Standard Shape / TextBox
            if shape.has_text_frame:
                if shape.text_frame.text.strip():
                    self._translate_text_frame(shape.text_frame, stats)

            # Table
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        tc = cell._tc
                        if tc not in self.visited_cells:
                            self.visited_cells.add(tc)
                            if cell.text_frame.text.strip():
                                self._translate_text_frame(cell.text_frame, stats)
                                stats["table_cells_processed"] += 1

            # Chart Title
            if shape.has_chart:
                chart = shape.chart
                if chart.has_title and chart.chart_title.has_text_frame:
                    self._translate_text_frame(chart.chart_title.text_frame, stats)

    def _translate_text_frame(self, text_frame, stats: Dict) -> None:
        """
        Translates a TextFrame paragraph by paragraph to preserve sentence grammar,
        reassembles runs, and applies complete OpenXML Times New Roman normalization.
        """
        has_any_cjk = False
        max_expansion = 1.0

        for paragraph in text_frame.paragraphs:
            raw_text = "".join(r.text for r in paragraph.runs) if paragraph.runs else paragraph.text
            if not raw_text.strip():
                continue

            if self.translator.contains_japanese(raw_text):
                has_any_cjk = True
                translated_text = self.translator.translate_text(raw_text)
                
                # Compute expansion ratio
                if len(raw_text) > 0:
                    expansion = len(translated_text) / len(raw_text)
                    max_expansion = max(max_expansion, expansion)

                # Reassemble runs: put translated text into first run, clear subsequent runs
                if paragraph.runs:
                    paragraph.runs[0].text = translated_text
                    for r in paragraph.runs[1:]:
                        r.text = ""
                else:
                    paragraph.text = translated_text

                stats["paragraphs_translated"] += 1

            # Enforce Times New Roman on the paragraph and all runs regardless of language
            apply_times_new_roman_complete(paragraph)
            for r in paragraph.runs:
                apply_times_new_roman_complete(r)

        # Normalize layout, margins, word wrap, and autofit
        OpenXMLTypographyNormalizer.normalize_text_frame(text_frame, expansion_ratio=max_expansion)
