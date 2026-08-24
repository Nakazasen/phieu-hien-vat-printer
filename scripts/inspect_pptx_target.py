import os
import sys
import json
import traceback
import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

FILES = {
    "file1": {
        "title": "Athena保証工程取り組み説明2025 VN.pptx",
        "path": r"\\10.170.162.32\Data\00_KDTVN Common(KDTVN共通)\⑤Production Engineering(製造技術)\◆Iris EXP◆\Điều chỉnh\Athenal\Athena保証工程取り組み説明2025 VN.pptx",
    },
    "file2": {
        "title": "Athena保証工程　RaspberryPI問題点 VN.pptx",
        "path": r"\\10.170.162.32\Data\00_KDTVN Common(KDTVN共通)\⑤Production Engineering(製造技術)\◆Iris EXP◆\Điều chỉnh\Athenal\Athena保証工程　RaspberryPI問題点 VN.pptx",
    }
}

def get_shape_type_str(st):
    if st is None:
        return "UNKNOWN"
    for attr in dir(MSO_SHAPE_TYPE):
        if not attr.startswith("_"):
            try:
                if getattr(MSO_SHAPE_TYPE, attr) == st:
                    return attr
            except Exception:
                pass
    return str(st)

def inspect_shape(shape, depth=0):
    st = getattr(shape, "shape_type", None)
    st_name = get_shape_type_str(st)
    
    shape_info = {
        "id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", ""),
        "type": st_name,
        "type_code": int(st) if st is not None else None,
        "left": getattr(shape, "left", 0),
        "top": getattr(shape, "top", 0),
        "width": getattr(shape, "width", 0),
        "height": getattr(shape, "height", 0),
        "depth": depth,
        "is_placeholder": False,
        "placeholder_type": None,
        "has_text": False,
        "text_content": "",
        "paragraphs": [],
        "is_table": False,
        "table_data": None,
        "is_picture": False,
        "image_info": None,
        "is_group": False,
        "children": [],
        "is_chart": False,
        "chart_title": None,
        "other_features": []
    }

    # Placeholder
    if getattr(shape, "is_placeholder", False):
        shape_info["is_placeholder"] = True
        try:
            shape_info["placeholder_type"] = str(shape.placeholder_format.type)
        except Exception:
            pass

    # Text frame
    if getattr(shape, "has_text_frame", False):
        shape_info["has_text"] = True
        paragraphs_data = []
        full_text = []
        try:
            for p in shape.text_frame.paragraphs:
                p_text = p.text
                if p_text:
                    full_text.append(p_text)
                runs_data = []
                for r in p.runs:
                    runs_data.append({
                        "text": r.text,
                        "font_name": getattr(r.font, "name", None),
                        "font_size_pt": r.font.size.pt if getattr(r.font, "size", None) else None,
                        "bold": getattr(r.font, "bold", None),
                        "italic": getattr(r.font, "italic", None),
                        "color": str(r.font.color.rgb) if getattr(r.font, "color", None) and hasattr(r.font.color, 'rgb') and r.font.color.rgb else None
                    })
                paragraphs_data.append({
                    "text": p_text,
                    "runs_count": len(p.runs),
                    "runs": runs_data
                })
            shape_info["text_content"] = "\n".join(full_text)
            shape_info["paragraphs"] = paragraphs_data
        except Exception as e:
            shape_info["text_error"] = str(e)

    # Table
    if getattr(shape, "has_table", False):
        shape_info["is_table"] = True
        try:
            table = shape.table
            grid = []
            for r_idx, row in enumerate(table.rows):
                row_data = []
                for c_idx, cell in enumerate(row.cells):
                    row_data.append({
                        "row": r_idx,
                        "col": c_idx,
                        "text": cell.text.strip(),
                        "paragraphs_count": len(cell.text_frame.paragraphs) if cell.text_frame else 0
                    })
                grid.append(row_data)
            shape_info["table_data"] = {
                "rows": len(table.rows),
                "cols": len(table.columns),
                "grid": grid
            }
        except Exception as e:
            shape_info["table_error"] = str(e)

    # Picture
    if st == getattr(MSO_SHAPE_TYPE, "PICTURE", None):
        shape_info["is_picture"] = True
        try:
            img = shape.image
            shape_info["image_info"] = {
                "content_type": getattr(img, "content_type", None),
                "ext": getattr(img, "ext", None),
                "size_bytes": len(img.blob) if hasattr(img, "blob") else 0,
                "dimensions": list(img.size) if hasattr(img, "size") else None, # (width, height) px
                "filename": getattr(img, "filename", None)
            }
        except Exception as e:
            shape_info["image_info"] = {"error": str(e)}

    # Chart
    if getattr(shape, "has_chart", False):
        shape_info["is_chart"] = True
        try:
            chart = shape.chart
            shape_info["chart_title"] = chart.chart_title.text_frame.text if chart.has_title else None
            shape_info["chart_type"] = str(chart.chart_type)
        except Exception as e:
            shape_info["chart_error"] = str(e)

    # Group
    if st == getattr(MSO_SHAPE_TYPE, "GROUP", None):
        shape_info["is_group"] = True
        try:
            for child in shape.shapes:
                shape_info["children"].append(inspect_shape(child, depth + 1))
        except Exception as e:
            shape_info["children_error"] = str(e)

    return shape_info

def analyze_presentation(file_key, file_meta):
    path = file_meta["path"]
    title = file_meta["title"]
    
    report = {
        "key": file_key,
        "title": title,
        "path": path,
        "exists": os.path.exists(path),
        "file_size": None,
        "readable": False,
        "writable": False,
        "slide_count": 0,
        "slide_dimensions": None,
        "aspect_ratio": None,
        "slides": [],
        "overall_stats": {
            "total_slides": 0,
            "slides_with_notes": 0,
            "total_shapes": 0,
            "shape_type_breakdown": {},
            "total_text_characters": 0,
            "total_tables": 0,
            "total_table_cells": 0,
            "total_pictures": 0,
            "total_image_bytes": 0,
            "image_formats": {},
            "total_group_shapes": 0,
            "total_charts": 0,
            "languages_detected": {
                "vietnamese_detected": False,
                "japanese_detected": False,
                "english_detected": False
            }
        },
        "errors": []
    }

    if not report["exists"]:
        report["errors"].append(f"Path does not exist: {path}")
        return report

    report["file_size"] = os.path.getsize(path)

    # Test read
    try:
        with open(path, "rb") as f:
            header = f.read(4096)
            report["readable"] = len(header) == 4096 or len(header) == report["file_size"]
    except Exception as e:
        report["errors"].append(f"Read error: {e}")

    # Test write safely: check if directory or file can be opened for append without altering content or check with os.access
    try:
        report["writable"] = os.access(path, os.W_OK)
        # also test directory write permission check
        parent_dir = os.path.dirname(path)
        report["dir_writable"] = os.access(parent_dir, os.W_OK)
    except Exception as e:
        report["errors"].append(f"Permission check error: {e}")

    # Open Presentation
    try:
        prs = Presentation(path)
        report["slide_count"] = len(prs.slides)
        report["overall_stats"]["total_slides"] = len(prs.slides)
        
        sw = prs.slide_width
        sh = prs.slide_height
        report["slide_dimensions"] = {
            "width_emu": sw,
            "height_emu": sh,
            "width_inches": round(sw / 914400, 2),
            "height_inches": round(sh / 914400, 2),
            "width_pt": round(sw / 12700, 2),
            "height_pt": round(sh / 12700, 2)
        }
        ratio = round(sw / sh, 3)
        if abs(ratio - 1.778) < 0.05:
            report["aspect_ratio"] = "16:9 (Widescreen)"
        elif abs(ratio - 1.333) < 0.05:
            report["aspect_ratio"] = "4:3 (Standard)"
        else:
            report["aspect_ratio"] = f"Custom ({ratio}:1)"

        all_text_blobs = []

        for s_idx, slide in enumerate(prs.slides):
            slide_entry = {
                "slide_number": s_idx + 1,
                "slide_id": slide.slide_id,
                "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
                "notes_text": "",
                "has_notes": False,
                "shape_count": len(slide.shapes),
                "shapes": [],
                "slide_stats": {
                    "text_boxes": 0,
                    "pictures": 0,
                    "tables": 0,
                    "groups": 0,
                    "charts": 0,
                    "shapes_by_type": {}
                }
            }

            # Check notes
            if slide.has_notes_slide:
                ntf = slide.notes_slide.notes_text_frame
                if ntf and ntf.text.strip():
                    slide_entry["notes_text"] = ntf.text.strip()
                    slide_entry["has_notes"] = True
                    report["overall_stats"]["slides_with_notes"] += 1
                    all_text_blobs.append(slide_entry["notes_text"])

            # Inspect shapes
            def process_shape_stats(s_info):
                t_str = s_info["type"]
                slide_entry["slide_stats"]["shapes_by_type"][t_str] = slide_entry["slide_stats"]["shapes_by_type"].get(t_str, 0) + 1
                report["overall_stats"]["shape_type_breakdown"][t_str] = report["overall_stats"]["shape_type_breakdown"].get(t_str, 0) + 1
                report["overall_stats"]["total_shapes"] += 1

                if s_info["has_text"] and s_info["text_content"]:
                    slide_entry["slide_stats"]["text_boxes"] += 1
                    report["overall_stats"]["total_text_characters"] += len(s_info["text_content"])
                    all_text_blobs.append(s_info["text_content"])

                if s_info["is_picture"] and s_info.get("image_info"):
                    slide_entry["slide_stats"]["pictures"] += 1
                    report["overall_stats"]["total_pictures"] += 1
                    sz = s_info["image_info"].get("size_bytes", 0)
                    report["overall_stats"]["total_image_bytes"] += sz
                    ext = s_info["image_info"].get("ext", "unknown")
                    report["overall_stats"]["image_formats"][ext] = report["overall_stats"]["image_formats"].get(ext, 0) + 1

                if s_info["is_table"] and s_info.get("table_data"):
                    slide_entry["slide_stats"]["tables"] += 1
                    report["overall_stats"]["total_tables"] += 1
                    report["overall_stats"]["total_table_cells"] += (s_info["table_data"]["rows"] * s_info["table_data"]["cols"])
                    for row in s_info["table_data"]["grid"]:
                        for cell in row:
                            if cell["text"]:
                                all_text_blobs.append(cell["text"])
                                report["overall_stats"]["total_text_characters"] += len(cell["text"])

                if s_info["is_group"]:
                    slide_entry["slide_stats"]["groups"] += 1
                    report["overall_stats"]["total_group_shapes"] += 1
                    for child in s_info.get("children", []):
                        process_shape_stats(child)

                if s_info.get("is_chart"):
                    slide_entry["slide_stats"]["charts"] += 1
                    report["overall_stats"]["total_charts"] += 1

            for shape in slide.shapes:
                s_info = inspect_shape(shape)
                slide_entry["shapes"].append(s_info)
                process_shape_stats(s_info)

            report["slides"].append(slide_entry)

        # Detect languages across combined text
        full_doc_text = "\n".join(all_text_blobs)
        jp_chars = sum(1 for c in full_doc_text if (0x3040 <= ord(c) <= 0x30FF) or (0x4E00 <= ord(c) <= 0x9FFF))
        vn_chars = sum(1 for c in full_doc_text if c in "àáảãạăằắẳẵặâầấẩẫậèéẻẽẹêềếểễệìíỉĩịòóỏõọôồốổỗộơờớởỡợùúủũụưừứửữựỳýỷỹỵđÀÁẢÃẠĂẰẮẲẴẶÂẦẤẨẪẬÈÉẺẼẸÊỀẾỂỄỆÌÍỈĨỊÒÓỎÕỌÔỒỐỔỖỘƠỜỚỞỠỢÙÚỦŨỤƯỪỨỬỮỰỲÝỶỸỴĐ")
        en_chars = sum(1 for c in full_doc_text if c.isascii() and c.isalpha())
        
        report["overall_stats"]["languages_detected"]["japanese_detected"] = jp_chars > 0
        report["overall_stats"]["languages_detected"]["japanese_char_count"] = jp_chars
        report["overall_stats"]["languages_detected"]["vietnamese_detected"] = vn_chars > 0
        report["overall_stats"]["languages_detected"]["vietnamese_char_count"] = vn_chars
        report["overall_stats"]["languages_detected"]["english_detected"] = en_chars > 0
        report["overall_stats"]["languages_detected"]["english_char_count"] = en_chars

    except Exception as e:
        report["errors"].append(f"PPTX parsing failure: {e}\n{traceback.format_exc()}")

    return report

def main():
    print("Beginning full extraction...")
    results = {}
    for k, meta in FILES.items():
        print(f"Processing {meta['title']}...")
        results[k] = analyze_presentation(k, meta)

    out_file = r"d:\Sandbox\PM_in_lai_phieuhienvat\scripts\pptx_deep_analysis.json"
    with open(out_file, "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    print(f"Deep analysis written to {out_file}")

if __name__ == "__main__":
    main()
