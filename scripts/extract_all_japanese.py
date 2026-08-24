import os
import sys
import pptx
import re
import json
from pptx.enum.shapes import MSO_SHAPE_TYPE

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TARGET_FILES = [
    os.path.join(PROJECT_ROOT, "backups", "pptx_inputs", "20260819_131424", "Athena保証工程取り組み説明2025 VN.pptx"),
    os.path.join(PROJECT_ROOT, "backups", "pptx_inputs", "20260819_131424", "Athena保証工程　RaspberryPI問題点 VN.pptx"),
]

CJK_REGEX = re.compile(r'[\u3040-\u30ff\u3400-\u4dbf\u4e00-\u9fff\uff66-\uff9f]')

all_jp_texts = set()

def inspect_text_frame(tf):
    for p in tf.paragraphs:
        full_text = "".join(r.text for r in p.runs) if p.runs else p.text
        if full_text and CJK_REGEX.search(full_text):
            all_jp_texts.add(full_text.strip())

def inspect_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            inspect_shapes(shape.shapes)
            continue
        if shape.has_text_frame:
            inspect_text_frame(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    inspect_text_frame(cell.text_frame)
        if shape.has_chart and shape.chart.has_title and shape.chart.chart_title.has_text_frame:
            inspect_text_frame(shape.chart.chart_title.text_frame)

for path in TARGET_FILES:
    if os.path.exists(path):
        prs = pptx.Presentation(path)
        for slide in prs.slides:
            inspect_shapes(slide.shapes)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                inspect_text_frame(slide.notes_slide.notes_text_frame)

print(f"Total unique Japanese paragraphs found: {len(all_jp_texts)}")
output_json = os.path.join(PROJECT_ROOT, "output", "extracted_japanese_texts.json")
with open(output_json, "w", encoding="utf-8") as f:
    json.dump(sorted(list(all_jp_texts)), f, ensure_ascii=False, indent=2)
print(f"Written to {output_json}")
